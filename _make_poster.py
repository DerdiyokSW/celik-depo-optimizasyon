"""GEÇİCİ — Poster taslağını PowerPoint (.pptx) olarak üretir. Bitince silinebilir.

Şartlar: 70×100 cm dikey, Times New Roman, ana başlık 70pt / konu başlığı 25pt / metin 20pt.
Gerçek ölçüm sayıları + TNR grafikleri (poster/fig*.png) + tablolar gömülü. Yazar/danışman
bilgileri PLACEHOLDER — kullanıcı doldurur. Çıktı: runs/evaluation/poster/POSTER.pptx
"""
import sys
from pathlib import Path

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt

ROOT = Path(".").resolve()
FIG = ROOT / "runs" / "evaluation" / "poster"
SIM = ROOT / "runs" / "evaluation" / "sim_state.png"
TNR = "Times New Roman"
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
GREEN = RGBColor(0x2E, 0x9E, 0x5B)
BLUE = RGBColor(0x3F, 0x7F, 0xC4)
GREY = RGBColor(0x55, 0x55, 0x55)

prs = Presentation()
prs.slide_width = Cm(70)
prs.slide_height = Cm(100)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # boş

# Sütun geometrisi (cm)
M = 2.5
GAP = 2.0
COLW = (70 - 2 * M - GAP) / 2  # 31.5
COLX = {0: M, 1: M + COLW + GAP}
cursor = {0: 15.0, 1: 15.0}


def _set_font(run, size, bold=False, color=None, italic=False):
    run.font.name = TNR
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def textbox(x, y, w, h, lines, size, bold=False, color=None, align=PP_ALIGN.LEFT,
            italic=False, anchor=MSO_ANCHOR.TOP):
    """lines: str ya da (str, dict) listesi. Çok satır destekler."""
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = ln
        _set_font(run, size, bold, color, italic)
    return tb


def add_heading(col, text):
    textbox(COLX[col], cursor[col], COLW, 1.6, text, 25, bold=True, color=NAVY)
    cursor[col] += 1.9


def add_body(col, text, height):
    textbox(COLX[col], cursor[col], COLW, height, text, 20)
    cursor[col] += height + 0.5


def add_caption(col, text):
    textbox(COLX[col], cursor[col], COLW, 0.9, text, 16, italic=True, color=GREY,
            align=PP_ALIGN.CENTER)
    cursor[col] += 1.1


def add_image(col, path, width=None):
    width = width or COLW
    w, h = Image.open(path).size
    height_cm = width * h / w
    x = COLX[col] + (COLW - width) / 2  # ortala
    slide.shapes.add_picture(str(path), Cm(x), Cm(cursor[col]), width=Cm(width))
    cursor[col] += height_cm + 0.3


def add_table(col, rows, col_widths, bold_row0=True, highlight=None):
    nr, nc = len(rows), len(rows[0])
    rowh = 1.25
    h = rowh * nr
    tbl = slide.shapes.add_table(nr, nc, Cm(COLX[col]), Cm(cursor[col]), Cm(COLW), Cm(h)).table
    for j, cw in enumerate(col_widths):
        tbl.columns[j].width = Cm(cw)
    for i, row in enumerate(rows):
        tbl.rows[i].height = Cm(rowh)
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            run = p.add_run(); run.text = str(val)
            is_hdr = (i == 0 and bold_row0)
            is_hi = highlight is not None and i == highlight
            _set_font(run, 18, bold=is_hdr or is_hi, color=RGBColor(255, 255, 255) if is_hdr else None)
            if is_hdr:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            elif is_hi:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xEA, 0xF6, 0xEE)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xF7, 0xF7, 0xF7)
    cursor[col] += h + 0.6


# ===================== BAŞLIK BLOĞU =====================
textbox(M, 1.5, 70 - 2 * M, 7.5,
        "Çelik Bobin Depo Yerleşiminde Rehandling ve Vinç Mesafesi Optimizasyonu",
        70, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
textbox(M, 9.5, 70 - 2 * M, 4.5,
        ["Sezgisel ve Pekiştirmeli Öğrenme (PPO) Yaklaşımlarının Karşılaştırması",
         "[Ad Soyad — Öğrenci No]   ·   [e-posta]",
         "Danışman: [Unvan Ad Soyad]   ·   Bilgisayar Mühendisliği Bölümü"],
        24, color=GREY, align=PP_ALIGN.CENTER)

# ===================== SOL SÜTUN =====================
add_heading(0, "Giriş ve Problem")
add_body(0, "Çelik fabrikalarında üretilen, her biri tonlarca ağırlığındaki bobinler önce bir "
            "depoya yerleştirilir, sonra kamyon, tren veya gemiye yüklenmek üzere alınır. Bobinler "
            "üst üste istiflenebildiği için, alınması gereken bir bobinin üstünde başka bir bobin "
            "varsa önce o kaldırılmalıdır — buna REHANDLING denir; gereksiz bir vinç hareketi, zaman "
            "ve enerji kaybıdır. Her hareket ayrıca bir VİNÇ YOLU (mesafe) doğurur. Bu proje, "
            "bobinleri en baştan akıllıca yerleştirerek bu iki maliyeti azaltmayı amaçlar.", 9.0)

add_heading(0, "Geliştirilen Sistem: Dijital İkiz")
add_body(0, ["Gerçek bir depoyu bilgisayarda canlandıran bir DİJİTAL İKİZ kurduk: 8 bölge × 36 sıra "
             "× 2 kat = 576 konum, tavan vinci, sürekli gelen bobinler ve değişen siparişler. Üzerine "
             "dört YERLEŞTİRME STRATEJİSİ takılıp canlı 3B panoda karşılaştırılabiliyor:",
             "•  Rastgele — kıyaslama tabanı",
             "•  Sezgisel — uzman kurallarıyla (aciliyet, kapıya yakınlık, istif düzeni)",
             "•  ML-Sezgisel — araç gecikmelerini tahmin eden makine öğrenmesi (LightGBM) ile",
             "•  PPO — deneme-yanılmayla KENDİ KENDİNE ÖĞRENEN pekiştirmeli öğrenme ajanı"], 11.5)

add_heading(0, "Geliştirilen 3B Pano")
add_image(0, FIG / "dash_main_ppo.png", width=29)
add_caption(0, "Şekil 1: Dijital ikizin canlı 3B görünümü. Her nokta bir bobin; renk aciliyeti "
               "gösterir (sarı = acil, mavi = bekleyebilir), yeşiller yükleme kapılarıdır.")

add_heading(0, "Yöntemleri Nasıl Değerlendirdik?")
add_body(0, "Dört yöntemi, daha önce HİÇ GÖRÜLMEMİŞ 30 depo senaryosunda, tıpatıp aynı koşullarda "
            "yarıştırdık (adil kıyaslama). Çalışma sırasında, öğrenen ajanın 'çok iyi' görünen iki "
            "sonucunun aslında yanıltıcı olduğunu fark edip düzelttik: önce EZBERLEME, sonra bazı "
            "kuralların GEVŞEMESİ. Bu öz-denetim, sonuçların güvenilir olmasını sağladı.", 8.0)
add_image(0, FIG / "fig3_metodoloji.png", width=29)
add_caption(0, "Şekil 2: Öğrenen ajanın gerçek başarısını bulmak için yapılan iki düzeltme.")

# ===================== SAĞ SÜTUN =====================
add_heading(1, "Bulgular: İki Çalışma Koşulu")
add_body(1, "Sistemi iki farklı koşulda inceledik; her birinde dört yöntem de aynı ortamda çalıştı.", 2.5)

add_body(1, "(A) GERÇEKÇİ DEPO — iki katlı istif ve her bobinin belirli bölgelere ait olduğu kurallar. "
            "Bu koşulda uzman-kuralı (sezgisel) yöntemler, en az rehandling ve en kısa vinç yolunu "
            "veren kararlı çözümlerdir.", 5.0)
add_table(1, [
    ["Yöntem", "Rehandling", "Vinç yolu (m)"],
    ["Rastgele", "33.1", "40 717"],
    ["ML-Sezgisel", "8.3", "29 806"],
    ["Sezgisel", "8.6", "26 525"],
    ["PPO (öğrenen)", "10.4", "35 675"],
], [13.5, 9.5, 11], highlight=3)
add_caption(1, "Çizelge 1: Gerçekçi depoda 30 senaryo ortalaması (düşük = iyi).")

add_body(1, "Öğrenen ajanın ilginç bir davranışı: depo AZ DOLUYKEN en az rehandling üreten yöntem oldu; "
            "depo dolduğunda ise hızla zorlandı. Yani gücü yük seviyesine bağlı.", 4.0)
add_image(1, FIG / "fig4_yuk_taramasi.png", width=31)
add_caption(1, "Şekil 3: Doluluk arttıkça yöntemlerin davranışı.")

add_body(1, "(B) SAF ROTA PROBLEMİ — tek katlı ve kısıtsız; tek amaç vinç yolunu kısaltmak. Bu, "
            "öğrenmenin doğal olarak güçlü olduğu koşuldur; öğrenen ajan burada uzman-kuralı yöntemden "
            "daha kısa vinç yolu buldu (30 senaryonun 27'sinde).", 5.0)
add_table(1, [
    ["Yöntem", "Vinç yolu (m)"],
    ["Rastgele", "39 621"],
    ["Sezgisel", "28 143"],
    ["PPO (öğrenen)", "26 650"],
], [16, 15], highlight=3)
add_caption(1, "Çizelge 2: Saf rota probleminde öğrenen ajan en kısa vinç yolunu bulur (p<0.001).")

add_heading(1, "Sonuç")
add_body(1, "Tek bir 'en iyi yöntem' yoktur; doğru yöntem PROBLEMİN YAPISINA göre değişir. Gerçekçi "
            "kısıtlı depoda uzman-kuralı yöntem güvenli ve kararlıdır; saf rota optimizasyonunda öğrenen "
            "ajan avantaj sağlar. Geliştirilen sistem her iki yöntemi de aynı arayüzde sunar ve canlı "
            "3B panoda izlenir — gerçek bir tesiste karar desteği için kullanılabilir.", 10.0)

OUT = FIG / "POSTER.pptx"
prs.save(str(OUT))
print("Kaydedildi ->", OUT)
print(f"Slayt boyutu: {prs.slide_width.cm:.0f}×{prs.slide_height.cm:.0f} cm")
print(f"Sol sütun son y: {cursor[0]:.1f} cm | Sağ sütun son y: {cursor[1]:.1f} cm (taşma için <98 olmalı)")
