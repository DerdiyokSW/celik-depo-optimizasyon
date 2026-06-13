"""GEÇİCİ — SUNUM-ICERIK.md'deki içerikten 16:9 SUNUM.pptx üretir (figürler gömülü)."""
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")

ROOT = Path(".").resolve()
POST = ROOT / "runs" / "evaluation" / "poster"
EVAL = ROOT / "runs" / "evaluation"
OUT = ROOT / "SUNUM.pptx"

NAVY = RGBColor(0x1F, 0x38, 0x64)
ACCENT = RGBColor(0xE0, 0x92, 0x2F)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF4, 0xF7, 0xFB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, size, color, bold, lvl) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.level = lvl
        p.space_after = Pt(6)
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
        r.font.name = "Calibri"
    return tb


def title_bar(slide, title, num):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tf = bar.text_frame; tf.margin_left = Inches(0.45); tf.margin_top = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = tf.paragraphs[0].add_run(); r.text = title
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.05), SW, Inches(0.06))
    acc.fill.solid(); acc.fill.fore_color.rgb = ACCENT; acc.line.fill.background()
    _txt(slide, SW - Inches(1.2), SH - Inches(0.5), Inches(1.0), Inches(0.4),
         [(f"{num}", 11, GREY, False, 0)], align=PP_ALIGN.RIGHT)
    _txt(slide, Inches(0.4), SH - Inches(0.5), Inches(8), Inches(0.4),
         [("Çelik Bobin Deposunda Yerleştirme Optimizasyonu — BSM 498", 10, GREY, False, 0)])


def bullets_runs(items):
    out = []
    for lvl, text in items:
        size = 19 if lvl == 0 else 15
        out.append((("• " if lvl == 0 else "– ") + text, size, DARK if lvl == 0 else GREY, False, lvl))
    return out


def add_image(slide, path, box):
    x, y, w, h = box
    img = Image.open(path); iw, ih = img.size
    scale = min(w / iw, h / ih)
    nw, nh = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(str(path), x + (w - nw) // 2, y + (h - nh) // 2, width=nw, height=nh)


def content_slide(title, items, num, image=None):
    s = prs.slides.add_slide(BLANK)
    title_bar(s, title, num)
    if image and Path(image).exists():
        _txt(s, Inches(0.5), Inches(1.4), Inches(6.3), Inches(5.6), bullets_runs(items),
             anchor=MSO_ANCHOR.TOP)
        add_image(s, image, (Inches(7.0), Inches(1.4), Inches(5.9), Inches(5.5)))
    else:
        _txt(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.5), bullets_runs(items))
    return s


# ---------- Slayt 1: Kapak ----------
s = prs.slides.add_slide(BLANK)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.0), SW, Inches(2.7))
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
_txt(s, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.4),
     [("Çelik Bobin Deposunda Yerleştirme Optimizasyonu", 34, WHITE, True, 0)], align=PP_ALIGN.CENTER,
     anchor=MSO_ANCHOR.MIDDLE)
_txt(s, Inches(0.8), Inches(3.55), Inches(11.7), Inches(1.0),
     [("Gecikme Tahmini ve Pekiştirmeli Öğrenme Destekli Karar Destek Sistemi", 18,
       RGBColor(0xDD, 0xE6, 0xF2), False, 0)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
_txt(s, Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.8),
     [("M. Yusuf Derdiyok — G231210381", 18, DARK, True, 0),
      ("Danışman: Dr. Öğr. Üyesi Hüseyin Demirci", 15, GREY, False, 0),
      ("Sakarya Üniversitesi · Bilgisayar Mühendisliği · BSM 498 · 2025–2026 Bahar", 13, GREY, False, 0)],
     align=PP_ALIGN.CENTER)
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.7), SW, Inches(0.08))
acc.fill.solid(); acc.fill.fore_color.rgb = ACCENT; acc.line.fill.background()

# ---------- İçerik slaytları ----------
SLIDES = [
    ("Problem", [
        (0, "Çelik bobinler (tonlarca) depoda istiflenir, sonra kamyon/tren/gemiye sevk edilir."),
        (0, "İki temel maliyet:"),
        (1, "Rehandling: alttaki bobini almak için üstündekini taşıma (üretken olmayan vinç hamlesi)."),
        (1, "Vinç mesafesi: her hareketin kat ettiği yol (enerji + çevrim süresi)."),
        (0, "Bu maliyetler büyük ölçüde YERLEŞTİRME kararına bağlıdır."),
        (0, "Karar belirsizlik altında: sevkiyat zamanı, gecikme, öncelik önceden tam bilinmez."),
        (0, "→ Hem kombinatoryal hem stokastik bir karar problemi."),
    ], None),
    ("Amaç ve Katkı", [
        (0, "Yerleştirme kararını destekleyen, uçtan uca ve modüler bir karar destek sistemi."),
        (0, "Dört stratejiyi AYNI ortamda kıyaslamak: Rastgele · Sezgisel · ML-Sezgisel · PPO."),
        (0, "Katkı 1 — Metodolojik titizlik: yalnızca görülmeyen (held-out) veride ölçüm + iki öz-düzeltme."),
        (0, "Katkı 2 — Yöntem-problem eşleşmesi: 'tek en iyi yöntem yok; problem yapısı belirler'."),
    ], None),
    ("Literatür (özet)", [
        (0, "SLAP (Storage Location Assignment): NP-zor; tam çözüm (MILP) yalnızca küçük örnekte."),
        (0, "Pratik: metaheuristikler (genetik, firework, yerel arama) ve son dönemde derin öğrenme."),
        (0, "PPO: kırpılmış vekil amaçla kararlı; farklı alanlarda (VizDoom, otonom sürüş, İHA) DQN'i geçer."),
        (0, "Boşluk: yöntemlerin aynı ortamda, görülmeyen veride, istatistiksel testle kıyaslanması sınırlı."),
    ], None),
    ("Sistem Mimarisi", [
        (0, "Dört katman: Veri Üretici → Simülasyon Çekirdeği ← Gecikme ML → Politikalar."),
        (0, "Hepsi arayüzsüz (headless); 3B görselleştirme üstte ince bir kabuk."),
        (0, "Ortak PlacementPolicy arayüzü → adil kıyas (aynı ortam) + risk izolasyonu."),
    ], POST / "fig_mimari.png"),
    ("Simülasyon Çekirdeği", [
        (0, "8 zone × 36 bay × 2 kat = 576 konum."),
        (0, "Tavan vinci köprü+araba EŞZAMANLI → Chebyshev yatay mesafe."),
        (0, "İstif sürekliliği + lojistik-hattı affinity (4 hat × 2 zone)."),
        (0, "single_layer / enforce_affinity bayrakları → iki yapılandırma tek çekirdekte."),
    ], POST / "fig_geometri.png"),
    ("Veri ve Olay Modeli", [
        (0, "Sentetik, deterministik: ~5000 bobin · 3600 araç · 1200 sipariş (seed → tekrar-üretilebilir)."),
        (0, "Poisson olay akışı: yeni sipariş 0,55 · gecikme 0,25 · iptal 0,10 · öncelik 0,10."),
        (0, "Gizli gecikme örüntüsü (mesafe × hava) → ML gürültü değil GERÇEK sinyal öğrenir."),
        (0, "Bilinçle-bozuk başlangıç (~%30, acil bobin gömülü) → iyileştirme gücünü ölçen zemin."),
    ], None),
    ("Gecikme Tahmin ML (hibrit kalp)", [
        (0, "LightGBM regresörü; bağımsız test MAE ≈ 6,95 dk."),
        (0, "Çıktı HEM ML-Sezgisel aciliyetine HEM PPO gözlemine beslenir → tahmin kararı doğrudan etkiler."),
        (0, "ML'in net katkısını izole eden tasarım (sezgisellerle aynı skor, tek fark gecikme)."),
    ], None),
    ("Yerleştirme Politikaları", [
        (0, "Rastgele (alt sınır) · Sezgisel (kural) · ML-Sezgisel (gecikme destekli) · PPO (öğrenen)."),
        (0, "Sezgisel ilke: acil bobin kapıya/erişilebilire, acil olmayan derine; affinity ikincil."),
        (0, "Skor (denklem): aciliyet u(c) (3.8) · erişilebilirlik a(slot) (3.9) · uyum (3.10)."),
    ], None),
    ("Yıldız Bileşen: PPO Ajanı", [
        (0, "MDP gözlemi: depo tensörü (8×36×2×3) + bekleyen bobin (gecikme dâhil) + küresel göstergeler."),
        (0, "Eylem: Discrete(576); eylem maskeleme geçersiz konumları eler (fizik ihlali öğrenilmez)."),
        (0, "Ödül: vinç maliyeti + rehandling cezası (yoğun); PPO clipped (3.5) + GAE (3.6)."),
        (0, "CNN özellik çıkarıcı; curriculum 8→12→20 olay/saat; 3.000.000 adım."),
    ], None),
    ("Değerlendirme Yöntemi", [
        (0, "Senaryo havuzu: 64 eğitim + 30 GÖRÜLMEYEN test (ayrık tohum, sızıntı yok)."),
        (0, "Eşleştirilmiş koşum: dört politika AYNI senaryo+olay tohumunu görür → fark yalnız politikadan."),
        (0, "Wilcoxon işaretli-sıra testi (eşleştirilmiş, dağılım-bağımsız); p < 0,05 anlamlı."),
    ], None),
    ("Metodolojik Titizlik (öne çıkan)", [
        (0, "Öz-düzeltme #1 (ezberleme): in-sample 3,40 → görülmeyen veride 21,62 → havuz ile 4,13."),
        (0, "Öz-düzeltme #2 (kısıt gevşemesi): 4,13 affinity kapalıyken yanıltıcıydı → affinity etkin 10,40."),
        (0, "Sonuçlar YALNIZCA held-out + gerçek kısıt altında raporlanır → güvenilirlik."),
    ], POST / "fig3_metodoloji.png"),
    ("Sonuç (1): Ana Senaryo", [
        (0, "2 kat + affinity, 30 held-out, 12 olay/saat."),
        (0, "En düşük: Sezgisel (rehandling 8,57 · vinç 26.525 m)."),
        (0, "PPO: rehandling 10,40 (p = 0,11, anlamsız), vinç 35.675 (p < 0,001, geride)."),
        (0, "→ Kısıtlı gerçekçi depoda sezgisel kararlı ve düşük maliyetli."),
    ], POST / "fig1_ana_senaryo.png"),
    ("Sonuç (2): Yük Bağımlılığı", [
        (0, "Düşük dolulukta (4/saat) PPO EN İYİ rehandling: 0,90."),
        (0, "Yük arttıkça sezgiseller öne; 20/saat'te PPO çöker (56,97 vs ~22)."),
        (0, "Vinç mesafesinde HER yükte sezgisel önde."),
        (0, "→ Rejime bağlı yöntem seçimi (düşük yük: öğrenen ajan; yüksek yük: sezgisel)."),
    ], POST / "fig4_yuk_taramasi.png"),
    ("Sonuç (3): Tek Katlı / Saf Rota", [
        (0, "İstif + affinity kalkar → saf uzamsal rota optimizasyonu (RL'in doğal alanı)."),
        (0, "PPO sezgiseli ANLAMLI geçer: 27/30 senaryo · %5,3 daha düşük vinç · p < 0,001."),
        (0, "150k adımda ≈ berabere → 3M adımda öne (ölçek etkisi)."),
    ], POST / "fig2_raf_senaryo.png"),
    ("Bulguların Yorumu", [
        (0, "İstif–yayılma ödünleşimi: PPO yerleşimi yayar → rehandling ↓ ama vinç yolu ↑."),
        (0, "Yöntem-problem eşleşmesi: tek 'en iyi' yok; istif/kısıt/yük yapısı doğru yöntemi belirler."),
        (0, "İki öz-düzeltmeyle birlikte: gözlenen fark ÖLÇÜM ARTEFAKTI DEĞİL, gerçek davranış farkı."),
    ], EVAL / "sim_state.png"),
    ("3B Dijital İkiz (canlı demo)", [
        (0, "Plotly/Dash canlı pano: bobinler aciliyet renkli; yan-yana iki politika kıyası."),
        (0, "Senaryo seçici: 2 kat ↔ tek kat canlı geçiş (PPO modeli otomatik swap)."),
        (0, "İsteğe bağlı canlı demo: python -m src.dashboard.app"),
    ], POST / "dash_main_ppo.png"),
    ("Mühendislik ve Ticarileşme", [
        (0, "Tekrar-üretilebilirlik: determinizm + train/test ayrımı + pytest (yedi katman)."),
        (0, "Donanımsız karar destek katmanı; hibrit dağıtım (düşük yük PPO, yüksek yük sezgisel)."),
        (0, "Düşük ilk yatırım; çıkarım milisaniye düzeyinde."),
    ], None),
    ("Sonuç ve Gelecek Çalışmalar", [
        (0, "Özet: 4 yöntem · held-out + Wilcoxon · iki öz-düzeltme · yöntem-problem eşleşmesi."),
        (0, "Gelecek: sevkiyatta vinç çizelgeleme (2. RL ajanı) · yüksek-yük curriculum ·"),
        (1, "MILP üst sınır referansı · gerçek tesis verisiyle doğrulama."),
    ], None),
]

for idx, (title, items, image) in enumerate(SLIDES, start=2):
    content_slide(title, items, idx, image)

# ---------- Son slayt: Teşekkür / Sorular ----------
s = prs.slides.add_slide(BLANK)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.6), SW, Inches(2.3))
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
_txt(s, Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.2),
     [("Teşekkürler — Sorular?", 36, WHITE, True, 0)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
_txt(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.4),
     [("Danışman: Dr. Öğr. Üyesi Hüseyin Demirci", 15, DARK, False, 0),
      ("Açık kaynak depo simülasyonu/SLAP katkısı: j4n1k", 13, GREY, False, 0),
      ("M. Yusuf Derdiyok — G231210381", 13, GREY, False, 0)], align=PP_ALIGN.CENTER)

prs.save(str(OUT))
print("Kaydedildi ->", OUT)
print("Slayt sayısı:", len(prs.slides._sldIdLst))
